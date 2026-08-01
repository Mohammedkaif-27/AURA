"""
AURA Backend — RAG System (Retrieval-Augmented Generation)

Purpose:
    Converts unstructured documents (PDF, DOCX) into searchable vectors,
    and retrieves relevant paragraphs to answer user questions.

Responsibilities:
    - Text Extraction: Read PDF/DOCX/PPTX, including OCR via Tesseract.
    - Chunking: Split text into small token-limited chunks.
    - Vector Database: Store and query embeddings using ChromaDB.
    - Hybrid Search: Combine Vector search + BM25 keyword search.
    - Reranking: Re-score search results for maximum accuracy.

Pipeline (Phases 3-5):
    ┌─────────────────────────────────────────────────────┐
    │  Query: "How to clean filter?"                      │
    │      │                                              │
    │      ▼                                              │
    │  Query Expansion (LLM)  →  generates 3-5 variants   │
    │      │                                              │
    │      ├───────────────────────┐                      │
    │      ▼                       ▼                      │
    │  Vector Search (Chroma)  BM25 Keyword Search        │
    │      │                       │                      │
    │      └───────────┬───────────┘                      │
    │                  ▼                                  │
    │      Reciprocal Rank Fusion (Merging)               │
    │                  │                                  │
    │                  ▼                                  │
    │      Cross-Encoder Reranking                        │
    │                  │                                  │
    │                  ▼                                  │
    │      Top K Context Chunks for Responder             │
    └─────────────────────────────────────────────────────┘

    Why API-based inference?
        We use HuggingFace Inference API (free tier) instead of loading
        models locally because it keeps memory under ~512 MB, enabling
        deployment on free-tier hosts (Render, Railway).

Used By:
    agents.py (retrieval_agent calls search_knowledge)

Depends On:
    chromadb, requests, rank_bm25, PyMuPDF, pytesseract
"""

import os
import gc
import glob
import time
import logging
import threading
import concurrent.futures
from datetime import datetime
from typing import List, Dict, Optional

import requests
import chromadb
from chromadb.config import Settings

from backend.config import (
    EMBEDDING_MODEL,
    RERANKER_MODEL,
    CHROMA_API_KEY,
    CHROMA_TENANT,
    CHROMA_DATABASE,
    CHROMA_COLLECTION,
    HF_API_KEY,
    AURA_VERSION,
    CHUNK_SIZE_TOKENS,
    CHUNK_OVERLAP_RATIO,
    RETRIEVAL_CANDIDATES,
    ENABLE_QUERY_EXPANSION,
    ENABLE_HYBRID_SEARCH,
)

logger = logging.getLogger(__name__)

# ── Module-level singletons ────────────────────────────────────────────────────────────────
_embedding_fn = None          # APIEmbeddingFunction instance
_embedding_dimension = None   # int — verified on startup
_chroma_client = None
_collection = None
_initialized = False

# Phase 5: BM25 index singletons
_bm25_index = None
_bm25_corpus: List[Dict] = []  # [{text, source, page, chunk_index, doc_id}, ...]

# ── Retrieval Cache ──
_retrieval_cache: Dict[str, List] = {}
_RETRIEVAL_CACHE_MAX_SIZE = 200
_cache_lock = threading.Lock()

# Thread safety locks
_init_lock = threading.Lock()
_bm25_lock = threading.Lock()


# ── API-based Models ──────────────────────────────────────────────────────────────────

from chromadb.api.types import EmbeddingFunction


class APIEmbeddingFunction(EmbeddingFunction):
    """ChromaDB-compatible embedding function using HF Inference API.

    Calls the feature-extraction endpoint for the configured EMBEDDING_MODEL.
    Includes retry/backoff for cold-start 503 and rate-limit 429 responses,
    matching the retry pattern used in llm_client.py.
    """

    _MAX_RETRIES = 4
    _BATCH_SIZE = 32  # Max texts per API call (conservative for serverless)

    def __init__(self):
        if not HF_API_KEY:
            raise RuntimeError(
                "HF_API_KEY is required for API-based embeddings. "
                "Get a free token at https://huggingface.co/settings/tokens"
            )
        self.api_url = f"https://router.huggingface.co/hf-inference/models/{EMBEDDING_MODEL}/pipeline/feature-extraction"
        self.headers = {"Authorization": f"Bearer {HF_API_KEY}"}

        # Warmup: verify API reachability and determine embedding dimension
        logger.info(f"Connecting to HF Inference API for embedding model: {EMBEDDING_MODEL}")
        warmup = self._embed_batch(["warmup"])
        if not warmup or not warmup[0]:
            raise RuntimeError(
                f"HF Inference API returned empty embeddings for {EMBEDDING_MODEL}. "
                "Check your HF_API_KEY and model availability."
            )
        self.dimension = len(warmup[0])

    def _embed_batch(self, texts: List[str]) -> List[List[float]]:
        """Embed a single batch of texts via HF Inference API with retry."""
        payload = {
            "inputs": texts,
            "options": {"wait_for_model": True},
        }

        for attempt in range(self._MAX_RETRIES):
            try:
                resp = requests.post(
                    self.api_url, headers=self.headers, json=payload, timeout=120
                )

                if resp.status_code in (503, 429):
                    wait = min(2 ** (attempt + 1), 30)
                    logger.warning(
                        f"HF Embedding API returned {resp.status_code}, "
                        f"retrying in {wait}s (attempt {attempt + 1}/{self._MAX_RETRIES})..."
                    )
                    time.sleep(wait)
                    continue

                resp.raise_for_status()
                result = resp.json()

                if not isinstance(result, list) or len(result) == 0:
                    raise ValueError(f"Unexpected API response shape: {type(result)}")

                # Handle 3D response (token-level embeddings) → mean pool to 2D
                if (
                    isinstance(result[0], list)
                    and len(result[0]) > 0
                    and isinstance(result[0][0], list)
                ):
                    import numpy as np
                    return [
                        np.mean(token_embs, axis=0).tolist()
                        for token_embs in result
                    ]

                # 2D response: already sentence-level embeddings
                return result

            except requests.exceptions.RequestException as e:
                if attempt == self._MAX_RETRIES - 1:
                    raise RuntimeError(
                        f"HF Embedding API failed after {self._MAX_RETRIES} retries: {e}"
                    )
                wait = min(2 ** (attempt + 1), 30)
                logger.warning(f"HF API request error (attempt {attempt + 1}): {e}")
                time.sleep(wait)

        raise RuntimeError("HF Embedding API failed — exhausted all retries")

    def name(self) -> str:
        return "hf_api_inference"

    def __call__(self, input: List[str]) -> List[List[float]]:
        # Batch large inputs to stay within API payload limits
        if len(input) <= self._BATCH_SIZE:
            return self._embed_batch(input)

        all_embeddings: List[List[float]] = []
        for i in range(0, len(input), self._BATCH_SIZE):
            batch = input[i : i + self._BATCH_SIZE]
            all_embeddings.extend(self._embed_batch(batch))
        return all_embeddings

    def embed_query(self, input: List[str]) -> List[List[float]]:
        return self.__call__(input)

    def embed_documents(self, input: List[str]) -> List[List[float]]:
        return self.__call__(input)


def _get_embedding_fn():
    """Return singleton API embedding function."""
    global _embedding_fn, _embedding_dimension
    if _embedding_fn is None:
        logger.info(f"Initializing API embedding function for: {EMBEDDING_MODEL}")
        start_time = time.time()
        try:
            _embedding_fn = APIEmbeddingFunction()
            _embedding_dimension = _embedding_fn.dimension
            elapsed = time.time() - start_time
            logger.info(
                f"Embedding API ready in {elapsed:.2f}s (dim: {_embedding_dimension})"
            )
        except (RuntimeError, requests.exceptions.RequestException) as e:
            logger.error(f"Failed to initialize embedding API for {EMBEDDING_MODEL}: {e}")
            raise
    return _embedding_fn


def _rerank_with_api(query: str, documents: List[str]) -> Optional[List[float]]:
    """Rerank documents using BAAI/bge-reranker-v2-m3 via HF Inference API.

    Uses batched REST payload to score all documents in a single HTTP request.
    Fails fast if unavailable to preserve graceful degradation.
    """
    if not HF_API_KEY or not documents:
        return None

    api_url = f"https://router.huggingface.co/hf-inference/models/{RERANKER_MODEL}"
    headers = {"Authorization": f"Bearer {HF_API_KEY}"}

    # Verified batched REST payload for HF text-classification
    payload = {
        "inputs": [{"text": query, "text_pair": doc} for doc in documents],
        "options": {"wait_for_model": True},
    }

    try:
        # 1 attempt, 5-second timeout to fail fast and prevent 210s latency
        resp = requests.post(api_url, headers=headers, json=payload, timeout=10)
        resp.raise_for_status()

        result = resp.json()
        logger.debug(f"RAW RERANKER RESPONSE (len={len(result) if isinstance(result, list) else 'N/A'}): {result}")
        
        # Response format: [[{"label": "LABEL_0", "score": 0.95}], [{"label": "LABEL_0", "score": 0.8}]]
        # Verified response shape: [[{...}, {...}, ..., {...}]] — one outer list
        # wrapping a single flat list of per-document scores.
        if isinstance(result, list) and len(result) == 1 and isinstance(result[0], list):
            result = result[0]

        scores = [item.get("score", 0.0) for item in result]

        while len(scores) < len(documents):
            scores.append(0.0)

        return scores[:len(documents)]

    except Exception as e:
        logger.warning(f"Batched Reranker API call failed: {e}")
        return None


# ── Self-Test & Validation ─────────────────────────────────────────────────────

def _run_startup_self_test():
    """Verify API endpoints and DB connections are functional before accepting requests."""
    logger.info("Running RAG startup self-test...")

    try:
        # 1. Check embedding API
        ef = _get_embedding_fn()
        if not ef:
            raise RuntimeError("Embedding function is None")

        # 2. Check reranker API (warm up — non-fatal if it fails)
        try:
            scores = _rerank_with_api("test query", ["test document"])
            if scores is not None:
                logger.info("Self-test: Reranker API is reachable")
            else:
                logger.warning(
                    "Self-test: Reranker API unavailable, system will run without it"
                )
        except Exception as e:
            logger.warning(f"Self-test: Reranker API check failed (non-fatal): {e}")

        # 3. Check ChromaDB collection
        global _collection
        if not _collection:
            raise RuntimeError("ChromaDB collection is None")

        doc_count = _collection.count()

        # 4. Check dimension compatibility
        if "embedding_dimension" in _collection.metadata:
            stored_dim = int(_collection.metadata["embedding_dimension"])
            if stored_dim != _embedding_dimension:
                raise RuntimeError(
                    f"Dimension mismatch! Collection was created with dim {stored_dim}, "
                    f"but {EMBEDDING_MODEL} produces dim {_embedding_dimension}. "
                    "You must rebuild the collection."
                )

        # 5. Dummy retrieval test
        if doc_count > 0:
            test_results = _collection.query(query_texts=["self test"], n_results=1)
            if not test_results:
                raise RuntimeError("ChromaDB retrieval returned invalid structure")
                
        logger.info(f"Self-test passed ✅ (Collection contains {doc_count} documents)")
        
    except Exception as e:
        logger.error(f"RAG Self-test failed ❌: {e}")
        raise RuntimeError(f"RAG startup validation failed: {e}")


# ── Token-aware recursive chunking ──────────────────────────────────

def _estimate_tokens(text: str) -> int:
    """Rough token estimate (~4 chars per token for English)."""
    return len(text) // 4


def _split_on_separators(text: str, separators: List[str]) -> List[str]:
    """Split text on the first matching separator, recursively."""
    if not separators:
        return [text]
    sep = separators[0]
    remaining_seps = separators[1:]
    parts = text.split(sep)
    results = []
    for part in parts:
        stripped = part.strip()
        if not stripped:
            continue
        if _estimate_tokens(stripped) <= CHUNK_SIZE_TOKENS:
            results.append(stripped)
        elif remaining_seps:
            results.extend(_split_on_separators(stripped, remaining_seps))
        else:
            # Hard split by character count as last resort
            char_limit = CHUNK_SIZE_TOKENS * 4
            for i in range(0, len(stripped), char_limit):
                results.append(stripped[i:i + char_limit].strip())
    return results


def chunk_text_recursive(text: str) -> List[str]:
    """
    Recursive token-based text splitter with overlap.
    Splits on paragraph → sentence → word boundaries in order.
    Returns chunks of ~CHUNK_SIZE_TOKENS tokens with ~15% overlap.
    """
    separators = ["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " "]
    raw_chunks = _split_on_separators(text, separators)

    # Merge small chunks and apply overlap
    overlap_tokens = int(CHUNK_SIZE_TOKENS * CHUNK_OVERLAP_RATIO)
    overlap_chars = overlap_tokens * 4

    merged: List[str] = []
    buffer = ""
    for chunk in raw_chunks:
        candidate = (buffer + " " + chunk).strip() if buffer else chunk
        if _estimate_tokens(candidate) <= CHUNK_SIZE_TOKENS:
            buffer = candidate
        else:
            if buffer:
                merged.append(buffer)
            buffer = chunk

    if buffer:
        merged.append(buffer)

    # Apply overlap: prepend tail of previous chunk to current
    if len(merged) <= 1:
        return merged

    overlapped: List[str] = [merged[0]]
    for i in range(1, len(merged)):
        prev_tail = merged[i - 1][-overlap_chars:] if len(merged[i - 1]) > overlap_chars else ""
        overlapped.append((prev_tail + " " + merged[i]).strip())

    return overlapped


# ── Initialization & Rebuild ─────────────────────────────────────────

def initialize_rag_system() -> bool:
    """Initialize ChromaDB, embedding model, reranker, and BM25 at startup.

    Thread-safe. All models are loaded before the first request.
    """
    global _chroma_client, _collection, _initialized

    with _init_lock:
        if _initialized:
            return True

        try:
            logger.info("Initializing RAG system (API-based embeddings)...")

            # 1. Initialize embedding API connection
            ef = _get_embedding_fn()

            # 2. Initialize ChromaDB Cloud Connection
            _chroma_client = chromadb.CloudClient(
                tenant=CHROMA_TENANT,
                database=CHROMA_DATABASE,
                api_key=CHROMA_API_KEY
            )

            # Metadata to verify compatibility later
            collection_metadata = {
                "hnsw:space": "cosine",
                "embedding_model": EMBEDDING_MODEL,
                "embedding_dimension": _embedding_dimension,
                "created_at": datetime.now().isoformat(),
                "aura_version": AURA_VERSION,
            }

            _collection = _chroma_client.get_or_create_collection(
                name=CHROMA_COLLECTION,
                embedding_function=ef,
                metadata=collection_metadata,
            )
            logger.info(f"ChromaDB initialized — collection '{CHROMA_COLLECTION}' "
                         f"({_collection.count()} existing docs)")

            # 3. Reranker — no preload needed (API calls are stateless)

            # 4. Phase 5: Build BM25 index from existing ChromaDB data
            if ENABLE_HYBRID_SEARCH:
                _rebuild_bm25_index()

            _initialized = True
            return True

        except Exception as e:
            logger.error(f"Failed to initialize RAG system: {e}", exc_info=True)
            return False


def rebuild_collection() -> bool:
    """
    Rebuild the Chroma collection by re-ingesting original source documents.
    Does NOT attempt to reuse existing Chroma data, as chunks/metadata may be incomplete.
    """
    global _collection, _chroma_client

    logger.warning("⚠️ REBUILDING CHROMA COLLECTION FROM SOURCE DOCUMENTS ⚠️")
    
    with _init_lock:
        try:
            if not _chroma_client:
                raise RuntimeError("Chroma client not initialized")
                
            # 1. Delete existing collection
            try:
                _chroma_client.delete_collection(name=CHROMA_COLLECTION)
                logger.info(f"Deleted old collection '{CHROMA_COLLECTION}'")
            except ValueError:
                pass # Collection doesn't exist yet
                
            # 2. Recreate with current embedding function and new metadata
            ef = _get_embedding_fn()
            
            collection_metadata = {
                "hnsw:space": "cosine",
                "embedding_model": EMBEDDING_MODEL,
                "embedding_dimension": _embedding_dimension,
                "created_at": datetime.now().isoformat(),
                "aura_version": AURA_VERSION,
            }
            
            _collection = _chroma_client.create_collection(
                name=CHROMA_COLLECTION,
                embedding_function=ef,
                metadata=collection_metadata,
            )
            logger.info(f"Created fresh collection '{CHROMA_COLLECTION}'")
            
        except Exception as e:
            logger.error(f"Failed to recreate collection: {e}")
            return False
            
    # 3. We don't trigger the Supabase fetch here directly because that's handled
    # in main.py via `_reindex_from_supabase()`. The caller should invoke that if needed.
        
    try:
        # 4. Rebuild BM25 index
        if ENABLE_HYBRID_SEARCH:
            _rebuild_bm25_index()
            
        logger.info(f"Collection rebuild complete. (Did not ingest local files - use Supabase).")
        return True
        
    except Exception as e:
        logger.error(f"Error during rebuilding index: {e}")
        return False


# ── Phase 5: BM25 Index Management ──────────────────────────────────

def _rebuild_bm25_index():
    """Rebuild the BM25 index from existing ChromaDB collection data."""
    global _bm25_index, _bm25_corpus

    with _bm25_lock:
        try:
            if not _collection:
                logger.warning("Cannot build BM25 index — collection not initialized")
                return

            count = _collection.count()
            if count == 0:
                logger.info("BM25 index: no documents to index")
                return

            # Fetch all documents from ChromaDB in chunks of 200 to avoid Cloud quota limits
            logger.info(f"Building BM25 index from {count} chunks...")
            all_ids, all_docs, all_metas = [], [], []
            offset = 0
            batch_size = 200
            
            while offset < count:
                batch_data = _collection.get(
                    include=["documents", "metadatas"],
                    limit=batch_size,
                    offset=offset,
                )
                if not batch_data or not batch_data.get("documents"):
                    break
                    
                all_ids.extend(batch_data["ids"])
                all_docs.extend(batch_data["documents"])
                all_metas.extend(batch_data["metadatas"])
                offset += batch_size

            if not all_docs:
                logger.warning("BM25 index: no documents returned from collection")
                return

            _bm25_corpus = []
            tokenized_corpus = []

            for doc_id, doc_text, meta in zip(all_ids, all_docs, all_metas):
                _bm25_corpus.append({
                    "text": doc_text,
                    "source": meta.get("source", "unknown"),
                    "page": meta.get("page", 0),
                    "chunk_index": meta.get("chunk_index", 0),
                    "doc_id": doc_id,
                })
                # Tokenize for BM25
                tokenized_corpus.append(doc_text.lower().split())

            if tokenized_corpus:
                from rank_bm25 import BM25Okapi
                _bm25_index = BM25Okapi(tokenized_corpus)
                logger.info(f"BM25 index built: {len(tokenized_corpus)} documents")
            else:
                logger.warning("BM25 index: empty corpus after tokenization")

        except Exception as e:
            logger.error(f"Failed to build BM25 index: {e}")


def _add_to_bm25_index(documents: List[str], metadatas: List[Dict], ids: List[str]):
    """Add new documents to the BM25 index (called during ingestion)."""
    global _bm25_index, _bm25_corpus

    if not ENABLE_HYBRID_SEARCH:
        return

    with _bm25_lock:
        try:
            for doc_text, meta, doc_id in zip(documents, metadatas, ids):
                _bm25_corpus.append({
                    "text": doc_text,
                    "source": meta.get("source", "unknown"),
                    "page": meta.get("page", 0),
                    "chunk_index": meta.get("chunk_index", 0),
                    "doc_id": doc_id,
                })

            # Rebuild the full BM25 index (BM25Okapi doesn't support incremental add)
            if _bm25_corpus:
                from rank_bm25 import BM25Okapi
                tokenized = [item["text"].lower().split() for item in _bm25_corpus]
                _bm25_index = BM25Okapi(tokenized)

        except Exception as e:
            logger.error(f"Failed to update BM25 index: {e}")


def _bm25_search(query: str, n: int = 20) -> List[Dict]:
    """Run BM25 keyword search and return top-n results with metadata."""
    if not _bm25_index or not _bm25_corpus:
        return []

    try:
        tokenized_query = query.lower().split()
        scores = _bm25_index.get_scores(tokenized_query)

        # Get top-n indices by score
        scored_indices = sorted(
            enumerate(scores), key=lambda x: x[1], reverse=True
        )[:n]

        results = []
        for idx, score in scored_indices:
            if score <= 0:
                continue
            item = _bm25_corpus[idx]
            results.append({
                "text": item["text"],
                "source": item["source"],
                "page": item["page"],
                "chunk_index": item["chunk_index"],
                "doc_id": item["doc_id"],
                "bm25_score": float(score),
            })

        return results

    except Exception as e:
        logger.error(f"BM25 search error: {e}")
        return []


def _reciprocal_rank_fusion(
    vector_results: List[Dict],
    bm25_results: List[Dict],
    k: int = 60,
    vector_weight: float = 1.0,
    bm25_weight: float = 1.5,
) -> List[Dict]:
    """Merge vector and BM25 results using weighted Reciprocal Rank Fusion.

    RRF score = weight * sum(1 / (k + rank_i)) across all result lists.
    BM25 is weighted slightly higher (1.5x) to prioritize exact keyword matches for technical manuals.
    Higher RRF score = better.
    """
    # Build a map: chunk_text -> {metadata, rrf_score}
    fused = {}

    for rank, item in enumerate(vector_results):
        key = item["text"][:200]  # Use first 200 chars as dedup key
        if key not in fused:
            fused[key] = {
                "text": item["text"],
                "source": item.get("source", "unknown"),
                "page": item.get("page", 0),
                "chunk_index": item.get("chunk_index", 0),
                "distance": item.get("distance", None),
                "rrf_score": 0.0,
            }
        fused[key]["rrf_score"] += vector_weight * (1.0 / (k + rank + 1))

    for rank, item in enumerate(bm25_results):
        key = item["text"][:200]
        if key not in fused:
            fused[key] = {
                "text": item["text"],
                "source": item.get("source", "unknown"),
                "page": item.get("page", 0),
                "chunk_index": item.get("chunk_index", 0),
                "distance": None,
                "rrf_score": 0.0,
            }
        fused[key]["rrf_score"] += bm25_weight * (1.0 / (k + rank + 1))

    # Sort by RRF score descending
    sorted_results = sorted(fused.values(), key=lambda x: x["rrf_score"], reverse=True)
    return sorted_results


# ── OCR via Tesseract (fast, CPU-native) ─────────────────────────────

_ocr_warning_logged = False  # Log the missing-dependency warning only once

def _ocr_page_tesseract(file_path: str, page_index: int) -> str:
    """OCR a single PDF page using PyMuPDF rendering + Tesseract.

    PyMuPDF renders the page to a 300-DPI pixmap (fast, even for
    vectorized-text PDFs like InDesign exports), then pytesseract
    runs Tesseract OCR on the resulting image.

    Returns extracted text, or empty string on failure.
    """
    global _ocr_warning_logged
    try:
        import fitz
        import pytesseract
        from PIL import Image
        import io as _io

        doc = fitz.open(file_path)
        page = doc[page_index]
        # 150 DPI balances OCR quality vs. speed/memory on constrained hosts
        pix = page.get_pixmap(dpi=150)
        img_bytes = pix.tobytes("png")
        del pix  # Free pixmap memory immediately
        doc.close()

        # Configure Tesseract path on Windows if not on PATH
        import shutil
        if not shutil.which("tesseract"):
            win_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            if os.path.isfile(win_path):
                pytesseract.pytesseract.tesseract_cmd = win_path
            else:
                if not _ocr_warning_logged:
                    logger.warning("Tesseract binary not found on PATH or default location")
                    _ocr_warning_logged = True
                return ""

        pil_image = Image.open(_io.BytesIO(img_bytes))
        del img_bytes  # Free raw PNG bytes
        text = pytesseract.image_to_string(pil_image)
        del pil_image  # Free PIL image
        return text.strip()

    except ImportError as ie:
        if not _ocr_warning_logged:
            logger.warning(f"OCR dependency missing: {ie}. "
                           "Install PyMuPDF and pytesseract (+ Tesseract binary).")
            _ocr_warning_logged = True
        return ""
    except Exception as e:
        logger.warning(f"Tesseract OCR failed for page {page_index + 1}: {e}")
        return ""


# ── Document extraction (multi-format) ───────────────────────────────

def _extract_pages_pdf(file_path: str, skip_ocr: bool = False, lightweight: bool = False) -> Dict:
    """Extract text and tables from a PDF.

    When lightweight=False (default): uses pdfplumber for high-quality table-aware
    extraction with optional OCR fallback.  ~300 MB RAM for large PDFs.

    When lightweight=True: uses pypdf for basic text extraction.  ~30 MB RAM.
    Used during startup re-indexing on memory-constrained hosts.

    Args:
        skip_ocr: If True, skip Tesseract OCR for pages without native text.
        lightweight: If True, use pypdf (low RAM) instead of pdfplumber.

    Returns a dict:
        {
            "pages": [{"page": int, "text": str}, ...],
            "total_pages": int,
            "ocr_pages": int,
            "ocr_required": bool,
        }
    """
    result = {"pages": [], "total_pages": 0, "ocr_pages": 0, "ocr_required": False}

    # ── Lightweight mode: pypdf (low memory) ─────────────────────────
    if lightweight:
        fallback_pages = _extract_pages_pdf_fallback(file_path)
        result["pages"] = fallback_pages
        result["total_pages"] = len(fallback_pages)
        logger.info(f"PDF extraction complete (lightweight/pypdf): "
                    f"{len(fallback_pages)} pages with text")
        return result

    # ── Full mode: pdfplumber (high quality, higher RAM) ─────────────
    try:
        import pdfplumber

        pages = []
        ocr_pages_count = 0
        native_text_pages = 0

        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            result["total_pages"] = total_pages

            for i, page in enumerate(pdf.pages):
                # extract_text(layout=True) preserves spatial formatting and tables!
                text = page.extract_text(layout=True) or ""

                if text and len(text.strip()) >= 50:
                    native_text_pages += 1
                elif not skip_ocr:
                    # OCR fallback via Tesseract (fast, CPU-native)
                    ocr_text = _ocr_page_tesseract(file_path, i)
                    if ocr_text and len(ocr_text.strip()) > 20:
                        text = ocr_text
                        ocr_pages_count += 1
                        logger.info(f"Page {i+1}: Tesseract OCR recovered "
                                    f"{len(ocr_text.strip())} chars")

                if text and text.strip():
                    pages.append({"page": i + 1, "text": text})

        # If zero pages had native text, the entire document is vectorized/scanned
        ocr_required = (native_text_pages == 0 and total_pages > 0)

        result.update({
            "pages": pages,
            "ocr_pages": ocr_pages_count,
            "ocr_required": ocr_required,
        })

        skipped_msg = " [OCR skipped — fast mode]" if skip_ocr else ""
        logger.info(f"PDF extraction complete: {len(pages)} pages with text "
                    f"({ocr_pages_count} via OCR) from {total_pages} total pages"
                    f"{' [FULL OCR — no native text layer]' if ocr_required else ''}"
                    f"{skipped_msg}")
        return result

    except ImportError:
        logger.error("pdfplumber not installed — pip install pdfplumber")
        fallback_pages = _extract_pages_pdf_fallback(file_path)
        result["pages"] = fallback_pages
        return result
    except Exception as e:
        logger.error(f"Error extracting PDF with pdfplumber: {e}")
        return result


def _extract_pages_pdf_fallback(file_path: str) -> List[Dict]:
    """Fallback PDF extraction using PyMuPDF first, then pypdf."""
    try:
        import fitz
        doc = fitz.open(file_path)
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text() or ""
            if text.strip():
                pages.append({"page": i + 1, "text": text})
        doc.close()
        return pages
    except ImportError:
        try:
            import logging
            # Silence the spammy KSCms-UHC-H font encoding errors from pypdf
            logging.getLogger("pypdf").setLevel(logging.CRITICAL)
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            pages = []
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    pages.append({"page": i + 1, "text": text})
            return pages
        except ImportError:
            logger.error("Neither PyMuPDF nor pypdf available for PDF extraction")
            return []
        except Exception as e:
            logger.error(f"Fallback PDF extraction failed with pypdf: {e}")
            return []
    except Exception as e:
        logger.error(f"Fallback PDF extraction failed with PyMuPDF: {e}")
        return []


def _extract_pages_docx(file_path: str) -> List[Dict]:
    """Extract text from a DOCX file, returning [{page, text}, ...].

    DOCX files don't have true 'pages', so we treat each ~2000-char block as a page.
    """
    try:
        from docx import Document
        doc = Document(file_path)
        full_text = "\n".join(para.text for para in doc.paragraphs if para.text.strip())

        if not full_text.strip():
            return []

        # Split into page-sized blocks
        block_size = 2000
        pages = []
        for i in range(0, len(full_text), block_size):
            block = full_text[i:i + block_size].strip()
            if block:
                pages.append({"page": (i // block_size) + 1, "text": block})
        return pages

    except ImportError:
        logger.error("python-docx not installed — pip install python-docx")
        return []
    except Exception as e:
        logger.error(f"Error extracting DOCX: {e}")
        return []


def _extract_pages_pptx(file_path: str) -> List[Dict]:
    """Extract text from a PPTX file, returning [{page, text}, ...].

    Each slide becomes one 'page'.
    """
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        pages = []

        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)
            slide_text = "\n".join(texts)
            if slide_text.strip():
                pages.append({"page": i + 1, "text": slide_text})

        return pages

    except ImportError:
        logger.error("python-pptx not installed — pip install python-pptx")
        return []
    except Exception as e:
        logger.error(f"Error extracting PPTX: {e}")
        return []


def _extract_pages_txt(file_path: str) -> List[Dict]:
    """Extract text from a plain .txt file, returning [{page, text}, ...].

    Splits into ~2000-char blocks, similar to DOCX handling.
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            full_text = f.read()

        if not full_text.strip():
            return []

        block_size = 2000
        pages = []
        for i in range(0, len(full_text), block_size):
            block = full_text[i:i + block_size].strip()
            if block:
                pages.append({"page": (i // block_size) + 1, "text": block})
        return pages

    except Exception as e:
        logger.error(f"Error extracting TXT: {e}")
        return []


def _extract_pages(file_path: str, skip_ocr: bool = False, lightweight: bool = False) -> Dict:
    """Route to the correct extractor based on file extension.

    Returns a dict with keys:
        - pages: List[Dict] — [{"page": int, "text": str}, ...]
        - total_pages: int
        - ocr_pages: int
        - ocr_required: bool
    Non-PDF formats return a simple wrapper with ocr_required=False.
    """
    ext = os.path.splitext(file_path)[1].lower()

    def _wrap(pages: List[Dict]) -> Dict:
        """Wrap non-PDF extractors into the standard metadata format."""
        return {
            "pages": pages,
            "total_pages": len(pages),
            "ocr_pages": 0,
            "ocr_required": False,
        }

    if ext == ".pdf":
        return _extract_pages_pdf(file_path, skip_ocr=skip_ocr, lightweight=lightweight)
    elif ext == ".docx":
        return _wrap(_extract_pages_docx(file_path))
    elif ext == ".pptx":
        return _wrap(_extract_pages_pptx(file_path))
    elif ext == ".txt":
        return _wrap(_extract_pages_txt(file_path))
    else:
        logger.warning(f"Unsupported file type: {ext}")
        return _wrap([])


# ── Document ingestion ───────────────────────────────────────────────

def ingest_single_document(file_path: str, source_name: str = "", doc_type: str = "manual", skip_ocr: bool = False, lightweight: bool = False) -> Dict:
    """
    Ingest a single document (PDF/DOCX/PPTX/TXT): extract → chunk → upsert into ChromaDB.

    Returns a dict:
        {
            "chunks_count": int,
            "ocr_pages": int,
            "total_pages": int,
            "ocr_required": bool,
        }

    Args:
        file_path: Path to the document file.
        source_name: Display name for citation metadata.
        doc_type: Document type tag — 'manual', 'policy', 'faq', or 'other'.
        skip_ocr: If True, skip Tesseract OCR (used during fast startup re-indexing).
        lightweight: If True, use pypdf instead of pdfplumber (low RAM mode).
    """
    if not _initialized:
        if not initialize_rag_system():
            raise RuntimeError("Cannot initialize RAG system")

    source = source_name or os.path.basename(file_path)
    extraction = _extract_pages(file_path, skip_ocr=skip_ocr, lightweight=lightweight)
    pages = extraction.get("pages", [])
    ocr_pages = extraction.get("ocr_pages", 0)
    total_pages = extraction.get("total_pages", 0)
    ocr_required = extraction.get("ocr_required", False)

    if not pages:
        logger.warning(f"No text extracted from {source}")
        return {"chunks_count": 0, "ocr_pages": ocr_pages,
                "total_pages": total_pages, "ocr_required": ocr_required}

    docs: List[str] = []
    ids: List[str] = []
    metadatas: List[Dict] = []
    chunk_idx = 0

    for page_info in pages:
        page_chunks = chunk_text_recursive(page_info["text"])
        for chunk in page_chunks:
            if not chunk.strip():
                continue
            docs.append(chunk)
            ids.append(f"{source}_p{page_info['page']}_c{chunk_idx}")
            metadatas.append({
                "source": source,
                "page": page_info["page"],
                "chunk_index": chunk_idx,
                "doc_type": doc_type,
            })
            chunk_idx += 1

    if not docs:
        return {"chunks_count": 0, "ocr_pages": ocr_pages,
                "total_pages": total_pages, "ocr_required": ocr_required}

    # ChromaDB handles embedding via the collection's embedding_function
    _collection.upsert(documents=docs, ids=ids, metadatas=metadatas)

    # Phase 5: Also add to BM25 index
    _add_to_bm25_index(docs, metadatas, ids)

    logger.info(f"Ingested {source}: {len(docs)} chunks (type={doc_type})")
    return {"chunks_count": len(docs), "ocr_pages": ocr_pages,
            "total_pages": total_pages, "ocr_required": ocr_required}



def ingest_policy_text(
    policy_id: str,
    text: str,
    policy_type: str = "",
    scope: str = "global",
    category: str = "",
) -> int:
    """Ingest a policy's description text into ChromaDB for RAG retrieval.

    Tags each chunk with policy-specific metadata so retrieval can filter
    by category when the user's query context includes a known product category.

    Returns the number of chunks created.
    """
    if not _initialized:
        if not initialize_rag_system():
            raise RuntimeError("Cannot initialize RAG system")

    if not text or not text.strip():
        logger.warning(f"Policy {policy_id}: empty description, skipping ingestion")
        return 0

    chunks = chunk_text_recursive(text)
    if not chunks:
        return 0

    docs: List[str] = []
    ids: List[str] = []
    metadatas: List[Dict] = []

    for i, chunk in enumerate(chunks):
        if not chunk.strip():
            continue
        chunk_id = f"policy_{policy_id}_c{i}"
        meta = {
            "source": f"policy:{policy_type}:{scope}",
            "page": 1,
            "chunk_index": i,
            "doc_type": "policy",
            "policy_type": policy_type,
            "policy_scope": scope,
        }
        if category:
            meta["category"] = category

        docs.append(chunk)
        ids.append(chunk_id)
        metadatas.append(meta)

    if not docs:
        return 0

    # Upsert into ChromaDB (overwrites previous version of same policy)
    _collection.upsert(documents=docs, ids=ids, metadatas=metadatas)

    # Also add to BM25 index
    _add_to_bm25_index(docs, metadatas, ids)

    logger.info(f"Policy {policy_id} ingested: {len(docs)} chunks (type={policy_type}, scope={scope}, category={category})")
    return len(docs)


# ── Phase 4: Query Expansion ────────────────────────────────────────

def _expand_query(query: str) -> List[str]:
    """Expand the user query into 3-5 semantic variants using the LLM.

    Phase 4: Helps bridge the gap between user phrasing and document terminology.
    Uses a fast, cheap LLM call with a tight prompt.
    """
    if not ENABLE_QUERY_EXPANSION:
        return [query]

    # Skip expansion for simple queries (e.g. fewer than 4 words)
    if len(query.strip().split()) < 4:
        return [query]

    try:
        from .llm_client import get_completion

        messages = [
            {
                "role": "system",
                "content": (
                    "You are a search query expansion assistant. "
                    "Given a user query, generate exactly 2 alternative phrasings that would help "
                    "find relevant information in product manuals, policy documents, and FAQs. "
                    "Include variations with:\n"
                    "- Technical terminology\n"
                    "- Common synonyms\n"
                    "Return ONLY a JSON array of strings, nothing else.\n"
                    "Example: [\"turbo wash feature\", \"turbo wash program mode\"]"
                ),
            },
            {
                "role": "user",
                "content": f"Expand this query: \"{query}\"",
            },
        ]

        resp = get_completion(messages, temperature=0.3, max_tokens=200)

        # Parse the JSON array from the response
        import json
        # Try to extract JSON array from response (handle markdown code blocks)
        resp_clean = resp.strip()
        if resp_clean.startswith("```"):
            resp_clean = resp_clean.split("\n", 1)[1] if "\n" in resp_clean else resp_clean
            resp_clean = resp_clean.rsplit("```", 1)[0].strip()

        expansions = json.loads(resp_clean)

        if isinstance(expansions, list) and all(isinstance(e, str) for e in expansions):
            # Always include the original query
            all_queries = [query] + [e for e in expansions if e.lower() != query.lower()]
            logger.info(f"Query expanded: '{query}' → {len(all_queries)} variants")
            return all_queries[:6]  # Cap at original + 5 expansions

    except Exception as e:
        logger.warning(f"Query expansion failed (using original only): {e}")

    return [query]


# ── Search / retrieval ───────────────────────────────────────────────

def search_knowledge(query: str, k: int = 5) -> List[str]:
    """
    Search the knowledge base.
    Returns a list of text chunks (top-k after optional reranking).
    """
    try:
        if not _initialized:
            logger.warning("RAG search skipped — system not initialized")
            return []

        doc_count = _collection.count() if _collection else 0
        if doc_count == 0:
            logger.warning("RAG search: collection is empty (0 documents indexed). "
                           "Upload documents via /admin/upload.")
            return []

        # Phase 3: Widen retrieval to RETRIEVAL_CANDIDATES (default 20)
        results = _collection.query(query_texts=[query], n_results=RETRIEVAL_CANDIDATES)

        if not results or not results["documents"] or not results["documents"][0]:
            logger.info(f"RAG search: no matches for query (collection has {doc_count} docs)")
            return []

        candidates = results["documents"][0]

        # Rerank via HF Inference API if available
        if len(candidates) > k:
            to_rerank = candidates[:10]
            scores = _rerank_with_api(query, to_rerank)
            if scores is not None:
                ranked = sorted(
                    zip(scores, to_rerank), key=lambda x: x[0], reverse=True
                )
                return [doc for _, doc in ranked[:k]]

        return candidates[:k]

    except Exception as e:
        logger.error(f"Error searching knowledge base: {e}")
        return []


def search_knowledge_with_metadata(query: str, k: int = 5, skip_expansion: bool = False) -> List[Dict]:
    """
    Search and return results WITH metadata (for citation chips in the UI).
    Returns: [{"text": ..., "source": ..., "page": ..., "chunk_index": ...}, ...]

    Pipeline (Phases 3-5):
      1. Expand query (Phase 4)
      2. Vector search for each variant (Phase 3: n=20)
      3. BM25 search for each variant (Phase 5)
      4. Merge via Reciprocal Rank Fusion (Phase 5)
      5. Rerank merged candidates (Phase 3)
      6. Return top-k
    """
    try:
        if not _initialized:
            logger.warning("RAG hybrid search skipped — system not initialized")
            return []

        doc_count = _collection.count() if _collection else 0
        if doc_count == 0:
            logger.warning("RAG hybrid search: collection is empty (0 documents indexed). "
                           "Upload documents via /admin/upload. Background indexing may still be in progress.")
            return []

        # ── Exact Query Cache ──
        cache_key = query.strip().lower()
        with _cache_lock:
            if cache_key in _retrieval_cache:
                logger.info(f"[CACHE HIT] RAG retrieval skipped for query: '{cache_key}'")
                return _retrieval_cache[cache_key][:k]

        # Phase 4: Expand query into variants
        if skip_expansion:
            query_variants = [query]
        else:
            query_variants = _expand_query(query)

        # Phase 3+5: Collect candidates from both vector and BM25 search
        all_vector_candidates = []
        all_bm25_candidates = []

        def _search_variant(variant: str):
            v_cands = []
            b_cands = []
            try:
                results = _collection.query(
                    query_texts=[variant],
                    n_results=RETRIEVAL_CANDIDATES,
                    include=["documents", "metadatas", "distances"],
                )
                if results and results["documents"] and results["documents"][0]:
                    for doc, meta, dist in zip(
                        results["documents"][0],
                        results["metadatas"][0],
                        results["distances"][0],
                    ):
                        v_cands.append({
                            "text": doc,
                            "source": meta.get("source", "unknown"),
                            "page": meta.get("page", 0),
                            "chunk_index": meta.get("chunk_index", 0),
                            "distance": dist,
                        })
            except Exception as e:
                logger.error(f"Vector search failed for variant '{variant}': {e}")

            if ENABLE_HYBRID_SEARCH:
                bm25_results = _bm25_search(variant, n=RETRIEVAL_CANDIDATES)
                b_cands.extend(bm25_results)
            return v_cands, b_cands

        # Execute variant searches in parallel to overlap network latency
        with concurrent.futures.ThreadPoolExecutor(max_workers=min(len(query_variants), 5)) as executor:
            futures = [executor.submit(_search_variant, v) for v in query_variants]
            for future in concurrent.futures.as_completed(futures):
                try:
                    v_cands, b_cands = future.result()
                    all_vector_candidates.extend(v_cands)
                    all_bm25_candidates.extend(b_cands)
                except Exception as e:
                    logger.error(f"Variant search thread failed: {e}")

        # Deduplicate vector candidates (keep best distance per chunk text)
        seen_texts = {}
        deduped_vector = []
        for c in all_vector_candidates:
            key = c["text"][:200]
            if key not in seen_texts or c["distance"] < seen_texts[key]["distance"]:
                seen_texts[key] = c
        deduped_vector = list(seen_texts.values())

        # Deduplicate BM25 candidates
        seen_bm25 = {}
        deduped_bm25 = []
        for c in all_bm25_candidates:
            key = c["text"][:200]
            if key not in seen_bm25 or c.get("bm25_score", 0) > seen_bm25[key].get("bm25_score", 0):
                seen_bm25[key] = c
        deduped_bm25 = list(seen_bm25.values())

        # Phase 5: Merge using Reciprocal Rank Fusion
        if ENABLE_HYBRID_SEARCH and deduped_bm25:
            # Sort vector by distance (ascending = best first)
            deduped_vector.sort(key=lambda x: x.get("distance", float("inf")))
            # Sort BM25 by score (descending = best first)
            deduped_bm25.sort(key=lambda x: x.get("bm25_score", 0), reverse=True)

            candidates = _reciprocal_rank_fusion(deduped_vector, deduped_bm25)
        else:
            # Vector-only: sort by distance
            deduped_vector.sort(key=lambda x: x.get("distance", float("inf")))
            candidates = deduped_vector

        if not candidates:
            return []

        # Phase 3: Rerank the top candidates via HF Inference API
        if len(candidates) > k:
            # Take top 10 for reranking (ensures < 2s latency while maintaining quality)
            to_rerank = candidates[:10]
            doc_texts = [c["text"] for c in to_rerank]
            scores = _rerank_with_api(query, doc_texts)
            if scores is not None:
                for c, s in zip(to_rerank, scores):
                    c["rerank_score"] = float(s)
                to_rerank.sort(key=lambda x: x.get("rerank_score", 0), reverse=True)
            candidates = to_rerank

        final_results = candidates[:k]
        
        # ── Store in Cache ──
        with _cache_lock:
            if len(_retrieval_cache) >= _RETRIEVAL_CACHE_MAX_SIZE:
                keys_to_evict = list(_retrieval_cache.keys())[: _RETRIEVAL_CACHE_MAX_SIZE // 10]
                for k_evict in keys_to_evict:
                    _retrieval_cache.pop(k_evict, None)
            _retrieval_cache[cache_key] = final_results

        return final_results

    except Exception as e:
        logger.error(f"Error searching knowledge base with metadata: {e}")
        return []
